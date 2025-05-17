from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import docx
import re
import tkinter as tk
from tkinter import filedialog

from datetime import datetime




# 設定樣式
title_font_size = Pt(60)
content_font_size = Pt(50)
page_font_size = Pt(18)
blue = RGBColor(0, 51, 102)
white = RGBColor(255, 255, 255)

# 建立簡報
ppt = Presentation()
layout = ppt.slide_layouts[6]

# 加頁函數
def add_formatted_slide(title, content, slide_number):
    slide = ppt.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = blue

    clean_title = re.sub(r"詩歌[:：].*$", "", title).strip()


    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean_title
    run.font.size = title_font_size
    run.font.bold = True
    run.font.color.rgb = white
    run.font.name = "微軟正黑體"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Pt(2) / 914400
    )
    line.fill.solid()
    line.fill.fore_color.rgb = white
    line.line.fill.background()

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.3))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True
    content_tf.clear()

    for line in content.split("\n"):
        if line.strip():
            p = content_tf.add_paragraph()
            run = p.add_run()
            run.text = line.strip()
            run.font.size = content_font_size
            run.font.color.rgb = white
            run.font.name = "微軟正黑體"
            p.alignment = PP_ALIGN.LEFT

    page_box = slide.shapes.add_textbox(Inches(9), Inches(6.8), Inches(1), Inches(0.5))
    page_tf = page_box.text_frame
    p = page_tf.paragraphs[0]
    run = p.add_run()
    run.text = str(slide_number)
    run.font.size = page_font_size
    run.font.color.rgb = white
    run.font.name = "微軟正黑體"
    p.alignment = PP_ALIGN.RIGHT

# 讀 Word 檔案
def parse_docx(path, title_lists):
    doc = docx.Document(path)
    current_title = None
    slides = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if any(title in text for title in title_lists):
            
            if "倫理一則" in text:
                current_title = "倫理一則"
                content = text.split("：", 1)[-1].strip() if "：" in text else text
                slides.append((current_title, content))
            else:
                current_title = text
        elif current_title:
            if not text.startswith("詩歌:"):
                slides.append((current_title, text))
            
    return slides


# -------------------original code-------------------

    # for para in doc.paragraphs:
    #     text = para.text.strip()
    #     print("原始文字：", text)
    #     print("______________________-")
    #     if not text:
    #         continue
    #     if text.startswith("一、") or text.startswith("二、") or text.startswith("三、") or \
    #        text.startswith("四、") or text.startswith("五、") or text.startswith("六、"):
    #         current_title = text
    #         print("標題：", current_title)
    #     elif text.endswith("之禱"):
    #         current_title = text
    #         print("標題：", current_title)
    #     elif text.startswith("倫理一則："):
    #         current_title = "倫理一則"
    #         content = text.split("：", 1)[-1].strip()
    #         slides.append((current_title, content))
    #     elif current_title:
    #         if not text.startswith("詩歌:"):
    #             slides.append((current_title, text))
                
    #             print("內容：", text)           
    #             print("-------------")
    #         else:
    #             print(" 詩歌：", text)

    # return slides


def extract_specific_text(text):
    """
    篩出 聖經片段
    
    """
    start_scripture = {"創", "出", "利", "民", "申", "書",
                            "士", "得", "撒上", "撒下", "王上", "王下", "代上", "代下",
                            "拉", "尼", "斯", "伯", "詩", "箴", "傳", "歌",
                            "賽", "耶", "哀", "結", "但", "何", "珥", "摩", "俄", "拿",
                            "彌", "鴻", "哈", "番", "該", "亞", "瑪",
                            "太", "可", "路", "約", "徒", "羅", "林", "林",
                            "加", "弗", "腓", "西", "帖", "帖", "提", "提",
                            "多", "門", "來", "雅", "彼前", "彼後", "約壹", "約一","約貳", "約二", "約三", "約參", "猶", "啟"}
    # 這裡可以根據需要調整正則表達式
    # 這裡的正則表達式是範例，實際情況可能需要根據文本格式進行調整
    pattern = r"(詩|數字):\s*\d+\s*~\s*\d+"
    matches = re.findall(pattern, text)
    return matches
    
# create the gui 
def create_gui_application(
):

    root = tk.Tk()
    # 固定視窗大小
    root.geometry("500x400")
    root.resizable(False, False)
    # 標題和版本
    root.title("禱告會PPT製作工具")
    version_label = tk.Label(root, text="Version: 1.0.1", fg="gray")
    version_label.pack(side=tk.BOTTOM, pady=5)
    input_docx_path = tk.Entry(root, width=40)
    # 選取檔案按鈕
    input_docx_path_label = tk.Label(root, text="請輸入或選取禱告單路徑")
    input_docx_path_label.pack(pady=5)
    select_file_button = tk.Button(root, text="選取", command=lambda: select_docx_file(input_docx_path))
    select_file_button.pack(pady=5)
    input_docx_path_label.pack(pady=5)
    input_docx_path.pack(pady=5)

    output_pptx_path_label = tk.Label(root, text="請輸入或選取輸出PPT名稱")
    output_pptx_path_label.pack(pady=5)
    output_pptx_path = tk.Entry(root, width=40)

    # 自動產生默認禱告會PPT output 名稱
    today_date = datetime.now().strftime("%Y年%m月%d日")
    output_pptx_path.insert(0, f"{today_date}_國度復興禱告會.pptx")
    output_pptx_path.pack(pady=5)

    # 這邊選取標題列表，請用dot 隔開(標題一,標題二)
    title_list_label = tk.Label(root, text="標題列表:")
    title_list_label.pack(pady=5)
    title_list_panel = tk.Entry(root, width=40)
    title_list_panel.pack(pady=5)

    # 產生標題列按鈕 
    execute_button = tk.Button(root, text="產生標題列表", command=lambda: get_spec_titles(input_docx_path,title_list_panel))

    execute_button.pack(pady=10)    
    # 產生PPT按鈕
    execute_button = tk.Button(root, text="產生PPT", command=lambda: generate_output_ppt(input_docx_path,output_pptx_path,title_list_panel))
    execute_button.pack(pady=10)
    



    root.mainloop()

def generate_output_ppt(input_docx_path, output_pptx_path, title_lists):



    docx_path =input_docx_path.get()
    if not docx_path:
        tk.messagebox.showerror("Error", "請選取文字檔")
        return
    elif not docx_path.endswith(".docx"):
        tk.messagebox.showerror("Error", "請選取正確的文字檔")
        return

    if title_lists.get() == "":
        tk.messagebox.showerror("Error", "請輸入標題列,或按產生標題列表")
        return  
    
    output_filename= output_pptx_path.get()
    if not output_filename:
        tk.messagebox.showerror("Error", "請輸入輸出的PPT名稱.")
        return
    
    title_lists = title_lists.get().split(",")
    slides_data = parse_docx(docx_path, title_lists)
    for i, (title, content) in enumerate(slides_data, 1):
        add_formatted_slide(title, content, i)
    
    
    ppt.save(output_filename)
    tk.messagebox.showinfo("Success", f"產生成功~路徑在{output_filename}!")

    

def get_spec_titles(input_docx_path, title_lists):
    """
        從 slides_data 中找出所有標題，回傳標題列表（去重、保留順序）
    """
    docx_path =input_docx_path.get()
    if not docx_path:
        tk.messagebox.showerror("Error", "請選取文字檔")
        return
    elif not docx_path.endswith(".docx"):
        tk.messagebox.showerror("Error", "請選取正確的文字檔")
        return
    
    doc = docx.Document(docx_path)
    
    title_list = []

    for para in doc.paragraphs:
        text = para.text.strip()
        
        if not text:
            continue
        if text.startswith("一、") or text.startswith("二、") or text.startswith("三、") or \
           text.startswith("四、") or text.startswith("五、") or text.startswith("六、"):
            clean_title = re.sub(r"詩歌[:：].*$", "", text).strip()
            title_list.append(clean_title)
        elif text.endswith("之禱"):
            clean_title = re.sub(r"詩歌[:：].*$", "", text).strip()
            title_list.append(clean_title)
        elif text.startswith("倫理一則：") or text.startswith("倫理一則"):
            title_list.append("倫理一則")

    
    # 去重並保留順序
    unique_titles = []
    seen = set()
    for t in title_list:
        if t not in seen:
            unique_titles.append(t)
            seen.add(t)
    title_lists.delete(0, tk.END)
    title_lists.insert(0, ",".join(unique_titles))

    # 將標題列表顯示在 GUI 上
    title_list_str = "\n".join(unique_titles)
             

    


    
    
    
    

def select_docx_file(input_docx_path):
    file_path = filedialog.askopenfilename(filetypes=[("Word Documents", "*.docx")])
    if file_path:
        input_docx_path.delete(0, tk.END)
        input_docx_path.insert(0, file_path)




if __name__ == "__main__":
    

    create_gui_application()
    
